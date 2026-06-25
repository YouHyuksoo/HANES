# 렌더된 slide-NN.png 들을 16:9 PPTX로 조립한다(슬라이드 1장 = 풀블리드 이미지 1장).
# 사용: python docs/presentation/scripts/build-pptx.py
import glob
import json
import os
from datetime import datetime, timezone, timedelta

from pptx import Presentation
from pptx.util import Inches

ROOT = os.getcwd()
BUILD_DIR = os.path.join(ROOT, "docs/presentation/.pptx-build")
OUT_PPTX = os.path.join(ROOT, "docs/presentation/hanes-mes-introduction.pptx")
MANIFEST = os.path.join(ROOT, "docs/presentation/artifact-build-manifest.json")
HTML = os.path.join(ROOT, "docs/presentation/hanes-mes-introduction.html")

EMU_W, EMU_H = Inches(13.333), Inches(7.5)  # 16:9

def main():
    pngs = sorted(glob.glob(os.path.join(BUILD_DIR, "slide-*.png")))
    if not pngs:
        raise SystemExit("렌더된 PNG가 없습니다. 먼저 render-slides.mjs를 실행하세요.")

    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    empty = 0
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=EMU_W, height=EMU_H)
        if os.path.getsize(png) < 5000:
            empty += 1
    prs.save(OUT_PPTX)

    kst = timezone(timedelta(hours=9))
    manifest = {
        "output": OUT_PPTX,
        "outputBytes": os.path.getsize(OUT_PPTX),
        "html": HTML,
        "htmlBytes": os.path.getsize(HTML),
        "slideCount": len(pngs),
        "slideSize": {"width": 1600, "height": 900, "ratio": "16:9"},
        "buildMode": "html-rendered-pptx",
        "packageCheck": {"mediaCount": len(pngs), "emptyMedia": empty},
        "updatedAt": datetime.now(kst).strftime("%Y-%m-%d %H:%M:%S KST"),
    }
    with open(MANIFEST, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    print(f"PPTX 생성: {len(pngs)}장 / {os.path.getsize(OUT_PPTX)} bytes / 빈이미지 {empty}")

if __name__ == "__main__":
    main()
